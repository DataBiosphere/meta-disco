"""Meta-Disco project-overview deck (AnVIL-branded). Built slide by slide."""

import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt

# anvilproject.org palette + typography
ANVIL_BLUE = RGBColor(0x03, 0x5C, 0x94)
ANVIL_TEAL = RGBColor(0x00, 0x72, 0x9C)
INK = RGBColor(0x21, 0x2B, 0x36)
GREEN = RGBColor(0x28, 0x75, 0x55)
GRAY = RGBColor(0x61, 0x61, 0x61)
FONT = "Inter"
MONO = "Roboto Mono"
CODE_BG = INK
CODE_FG = RGBColor(0xE6, 0xED, 0xF3)
CODE_ACCENT = RGBColor(0x7E, 0xC8, 0xE3)

LOGO = os.path.join(os.path.dirname(os.path.abspath(__file__)), "logoAnvil_tiny.png")

PART1_JSON = """{
 "file_name": "NA12485.cram",
 "md5sum": "0834b7799d163d48c08fd5243296079a",
 "file_size": 11227299324,
 "file_format": ".cram",
 "dataset_title": "ANVIL_T2T",
 "classifications": {
  "data_modality": {
   "value": "genomic",
   "status": "classified",
   "evidence": [
    {
     "rule_id": "program_bwa",
     "reason": "BWA (Burrows-Wheeler Aligner) is the standard short-read aligner for DNA sequencing. Commonly used for WGS, WES, and ChIP-seq",
     "tier": 3,
     "value": "genomic"
    }
   ]
  },
  "data_type": {
   "value": "alignments",
   "status": "classified",
   "evidence": [
    {
     "rule_id": "program_bwa",
     "reason": "BWA (Burrows-Wheeler Aligner) is the standard short-read aligner for DNA sequencing. Commonly used for WGS, WES, and ChIP-seq",
     "tier": 3,
     "value": "alignments"
    }
   ]
  },"""

PART2_JSON = """  "platform": {
   "value": "ILLUMINA",
   "status": "classified",
   "evidence": [
    {
     "rule_id": "platform_illumina",
     "reason": "PL:ILLUMINA in @RG header \\u2014 definitive Illumina platform tag",
     "tier": 3,
     "value": "ILLUMINA"
    }
   ]
  },
  "reference_assembly": {
   "value": "CHM13",
   "status": "classified",
   "evidence": [
    {
     "rule_id": "contig_length_detection",
     "reason": "Reference CHM13 detected from 18 matching contig lengths (definitive)",
     "tier": 4,
     "value": "CHM13"
    }
   ]
  },
  "assay_type": {
   "value": "WGS",
   "status": "classified",
   "evidence": [
    {
     "rule_id": "infer_assay_type",
     "reason": "Inferred WGS from platform/modality/file size signals",
     "tier": 3,
     "value": "WGS"
    }
   ]
  }
 },
 "entry_id": "4e466c25-a99f-4bbf-b944-3d9678ae80b1"
}"""


prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)


def add_logo(slide):
    """Same size, same corner, every slide."""
    slide.shapes.add_picture(LOGO, Inches(11.5), Inches(6.65), width=Inches(1.6))


def color_title(slide, rgb):
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = rgb


def brand_fonts(prs):
    """Set the brand font on every non-mono run (no theme-level hook in pptx)."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name != MONO:
                        run.font.name = FONT


def add_code_block(slide, code: str, left, top, width, height, size=14, accent_prefixes=(), spacing=1.15, wrap=False):
    """Dark panel with monospace text."""
    box = slide.shapes.add_textbox(left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Inches(0.25)
    tf.margin_top = tf.margin_bottom = Inches(0.18)
    first = True
    for line in code.splitlines():
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = para.add_run()
        run.text = line if line else " "
        run.font.name = MONO
        run.font.size = Pt(size)
        para.line_spacing = spacing
        accent = any(line.lstrip().startswith(p) for p in accent_prefixes)
        run.font.color.rgb = CODE_ACCENT if accent else CODE_FG
    return box


def bullets(slide, items, size=20):
    """Fill the body placeholder: items are (text, level, color-or-None)."""
    body = slide.placeholders[1].text_frame
    first = True
    for text, level, rgb in items:
        para = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        para.text = text
        para.level = level
        para.font.size = Pt(size)
        para.font.color.rgb = rgb or INK
        para.space_after = Pt(16)
        para.line_spacing = 1.15
    return body



def term_bullets(slide, items, size=22):
    """Bullets of the form '<bold term> — rest' (two runs per paragraph)."""
    body = slide.placeholders[1].text_frame
    first = True
    for term, rest in items:
        para = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        r1 = para.add_run()
        r1.text = term
        r1.font.bold = True
        r1.font.size = Pt(size)
        r1.font.color.rgb = INK
        r2 = para.add_run()
        r2.text = " — " + rest
        r2.font.size = Pt(size)
        r2.font.color.rgb = INK
        para.space_after = Pt(16)
        para.line_spacing = 1.15
    return body



ISSUE_URL = "https://github.com/DataBiosphere/meta-disco/issues/{n}"
_ISSUE_REF = __import__("re").compile(r"#(\d+)")


def linkify_issue_refs(prs):
    """Turn every #NNN in any text run into a hyperlink to the GitHub issue."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                runs = list(para.runs)
                if not any(_ISSUE_REF.search(r.text or "") for r in runs):
                    continue
                # capture (text, style, link) segments, then rebuild the runs
                segments = []
                for r in runs:
                    style = (r.font.name, r.font.size, r.font.bold, r.font.color.rgb if r.font.color and r.font.color.type is not None else None)
                    text = r.text or ""
                    pos = 0
                    for m in _ISSUE_REF.finditer(text):
                        if m.start() > pos:
                            segments.append((text[pos : m.start()], style, None))
                        segments.append((m.group(0), style, ISSUE_URL.format(n=m.group(1))))
                        pos = m.end()
                    if pos < len(text):
                        segments.append((text[pos:], style, None))
                for r in runs:
                    r._r.getparent().remove(r._r)
                for text, (fname, fsize, fbold, fcolor), link in segments:
                    nr = para.add_run()
                    nr.text = text
                    if fname:
                        nr.font.name = fname
                    if fsize:
                        nr.font.size = fsize
                    if fbold is not None:
                        nr.font.bold = fbold
                    if fcolor is not None:
                        nr.font.color.rgb = fcolor
                    if link:
                        nr.hyperlink.address = link


# ---------------------------------------------------------------- Slide 1: title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Meta-Disco (MetaData Discovery)"
slide.placeholders[1].text = "Project overview — current state and next steps\nAugust 18, 2026"
color_title(slide, ANVIL_BLUE)
for para in slide.placeholders[1].text_frame.paragraphs:
    para.font.color.rgb = GRAY
    para.font.size = Pt(22)
add_logo(slide)

# ------------------------------------------------- Slide 2: what meta-disco is
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "What Meta-Disco is (a reminder)"
color_title(slide, ANVIL_BLUE)
bullets(
    slide,
    [
        ("Discovers metadata for omics data files — BAM/CRAM, VCF, FASTQ, FASTA, GFA, BED, …", 0, None),
        ("LLM authors rules that are run by a deterministic rule engine", 0, None),
        ("No LLM at inference time: deterministic Python only — reproducible, re-runnable, cheap", 0, None),
        ("Every value carries full provenance: which rule fired and why", 0, None),
    ],
    size=22,
)
add_logo(slide)

# ------------------------------------------------- Slide 3: the five dimensions
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Current metadata dimensions"
color_title(slide, ANVIL_TEAL)
term_bullets(
    slide,
    [
        ("data_modality", "what the data measures (genomic, transcriptomic.*, epigenomic.*, imaging.histology)"),
        ("data_type", "what the file contains (alignments, variant calls, sequence, assembly, pangenome, …)"),
        ("reference_assembly", "the coordinate system (GRCh37, GRCh38, CHM13)"),
        ("assay_type", "the experiment that produced it (WGS, WES, RNA-seq, …)"),
        ("platform", "the instrument (ILLUMINA, PACBIO, ONT, …)"),
    ],
    size=22,
)
add_logo(slide)

# ------------------------------------- Slide 4: example rule (simplest case)
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Example Rule — tier 1 (extension)"
color_title(slide, ANVIL_BLUE)
add_code_block(
    slide,
    """# unified_rules.yaml — the simplest kind of rule (complete, verbatim)
- id: variant_default_genomic
  tier: 1
  scope: extension
  when:
    extensions: [".vcf", ".bcf", ".gvcf", ".g.vcf"]
  then:
    data_modality: genomic
    data_type: variants
  rationale: "VCF files contain variant calls (genomic data)\"""",
    Inches(0.9),
    Inches(1.7),
    Inches(11.5),
    Inches(4.4),
    size=17,
    accent_prefixes=("- id:", "tier:"),
)
add_logo(slide)

# ---------------------------------------------------- Slide 5: example rule
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Example Rule — tier 3 (header)"
color_title(slide, ANVIL_BLUE)
add_code_block(
    slide,
    """# unified_rules.yaml — one of the LLM-authored rules (complete, verbatim)
- id: platform_illumina
  tier: 3            # tier = priority, higher wins: 4 content-read > 3 header > 2 filename > 1 extension
  scope: header      # evaluated only when the file's header was fetched
  when:              # ALL conditions must hold for the rule to fire
    extensions: [".bam", ".cram"]
    header_section: "@RG"
    header_field:   "PL"
    header_pattern: "ILLUMINA"
  then:
    platform: ILLUMINA
  rationale: "PL:ILLUMINA in @RG header — definitive Illumina platform tag\"""",
    Inches(0.9),
    Inches(1.7),
    Inches(11.5),
    Inches(4.6),
    size=17,
    accent_prefixes=("- id:", "tier:"),
)
add_logo(slide)

# --------------------- Slides 6-7: example output, full record JSON (split)
for idx, part in enumerate([PART1_JSON, PART2_JSON], start=1):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Example Output — the full record ({idx}/2)"
    color_title(slide, ANVIL_BLUE)
    add_code_block(
        slide,
        part,
        Inches(0.7),
        Inches(1.45),
        Inches(12.0),
        Inches(5.55),
        size=10,
        spacing=1.0,
        wrap=True,
        accent_prefixes=('"data_modality"', '"data_type"', '"platform"', '"reference_assembly"', '"assay_type"', '"file_name"'),
    )
    add_logo(slide)

# ----------------------------- Slide 9: not classifying is an explicit outcome
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "Not classifying is an explicit outcome"
color_title(slide, ANVIL_BLUE)
box = slide.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(11.5), Inches(1.9))
tf = box.text_frame
tf.word_wrap = True
for i, (term, rest) in enumerate([
    ("classified", "evidence supports a value"),
    ("not_applicable", "the dimension makes no sense for this file kind — an index has no platform; a de novo assembly has no reference_assembly (it is its own coordinate system)"),
    ("not_classified", "applicable, but evidence is insufficient — better empty than guessed"),
    ("conflict", "top claims tie — flagged for a human, never silently resolved"),
]):
    para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    r1 = para.add_run(); r1.text = term; r1.font.bold = True; r1.font.size = Pt(17); r1.font.color.rgb = INK
    r2 = para.add_run(); r2.text = " — " + rest; r2.font.size = Pt(17); r2.font.color.rgb = INK
    para.space_after = Pt(6)
add_code_block(
    slide,
    """# rules can declare statuses, not just values (complete, verbatim)
- id: index_not_applicable
  tier: 1
  scope: extension
  when:
    extensions: [".bai", ".crai", ".tbi", ".csi", ".fai", ".idx", ".pbi"]
  then:
    status:
      data_modality: not_applicable
      data_type:     not_applicable
      assay_type:    not_applicable
      platform:      not_applicable
  rationale: "Index file — metadata inherited from parent data file\"""",
    Inches(0.9),
    Inches(3.35),
    Inches(11.5),
    Inches(3.15),
    size=13,
    spacing=1.05,
    accent_prefixes=("- id:", "status:"),
)
add_logo(slide)

# ---------------------------------------- Slide 10: validation framing
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "How do we know it's right?"
color_title(slide, ANVIL_TEAL)
term_bullets(
    slide,
    [
        ("External truth", "compare our answers, file by file, against independent authorities that published their own metadata"),
        ("Internal consistency", "inputs and outputs match their declared shapes (allowed values, required fields) and never contradict each other"),
        ("Three outcomes tracked", "agree · disagree · unable to answer — disagreements and gaps are the interesting part"),
    ],
    size=22,
)
add_logo(slide)

# ------------------------------------- Slide 11: external truth 1/3 — HPRC
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "External truth — HPRC Catalog"
color_title(slide, ANVIL_TEAL)
term_bullets(
    slide,
    [
        ("The universe", "we classify the full HPRC Catalog directly — 15,436 records (sequencing 6,048 · alignments 89 · annotations 8,739 · assemblies 560); 9,182 get classification records, the rest are formats without rules yet (tsv/txt/bigwig/gff3/indexes)"),
        ("platform", "6,048 of 6,048 sequencing files agree · 0 disagree · 0 unanswered"),
        ("modality", "of the same 6,048: 929 agree · 0 disagree · 4,919 unanswered (raw-reads ceiling) · 200 with no Catalog value to compare"),
        ("assay", "of the same 6,048: 929 agree · 0 disagree · 5,119 unanswered"),
        ("reference", "2,574 files carry a declared reference (alignments + annotations): 2,563 agree · 4 disagree (pangenome graphs, no single right answer — #331) · 7 unanswered"),
        ("Bridge to AnVIL", "~90% of the HPRC Catalog's files are also in AnVIL — a good testing example"),
    ],
    size=16,
)
add_logo(slide)

# ---------------------------- Slide 12: external truth 2/3 — 1000G and ENA
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "External truth — 1000 Genomes & ENA"
color_title(slide, ANVIL_TEAL)
term_bullets(
    slide,
    [
        ("1000 Genomes (IGSR)", "3,208 of 3,211 samples resolved at IGSR (79 lookup errors) → 14,780 files compared; IGSR lists each sample's platforms as a set, so agree = our value is in the set"),
        ("platform", "12,877 agree · 6 disagree · 1,897 with no committed value on our side — sums to 14,780"),
        ("modality", "14,160 agree (2,351 of those at family level, e.g. transcriptomic vs transcriptomic.bulk) · 0 disagree · 620 with no committed value — sums to 14,780"),
        ("The 6 disagreements", "our headers say ONT; the IGSR collections for those samples list only PacBio/Illumina — flagged for investigation, each a real lead"),
        ("ENA", "6,962 accession-bearing FASTQs vs the European Nucleotide Archive (re-run 2026-08-18): platform 6,830/6,830 agree, 100% (132 lookup errors); ENA also declares library source/strategy — the very fields FASTQs cannot reveal — an import source, not just a check"),
    ],
    size=16,
)
add_logo(slide)

# --------------- Slide 13: external truth 3/3 — AnVIL itself, Ensembl/NCBI
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "External truth — AnVIL itself & Ensembl/NCBI"
color_title(slide, ANVIL_TEAL)
term_bullets(
    slide,
    [
        ("AnVIL's own metadata", "per-file data_modality / reference_assembly compared where present"),
        ("Ensembl / NCBI", "a different kind of check: our contig-length tables are validated against the official assembly definitions"),
    ],
    size=19,
)
add_logo(slide)

# ------------------------------------ Slide 13: internal consistency, plainer
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Internal consistency: the pipeline checks itself"
color_title(slide, ANVIL_TEAL)
term_bullets(
    slide,
    [
        ("Before a run", "the downloaded input metadata must match the expected shape, or we refuse to start"),
        ("After a run", "every record is checked against the formal schema — controlled vocabulary only, status always present"),
        ("Cross-field linter", "logic rules over the whole corpus (an index file must stay not_applicable) — currently zero violations"),
        ("Determinism", "same inputs, same outputs — so every change is re-run corpus-wide and diffed"),
    ],
    size=19,
)
add_logo(slide)

# ------------------------------------ Slide 12: where we cannot answer, and why
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Where we can't answer — the content ceiling"
color_title(slide, ANVIL_BLUE)
term_bullets(
    slide,
    [
        ("Coverage today", "share of ~669K data files with a committed value (the 734K snapshot minus non-data files; not_applicable does not count): data_type ~97% · data_modality ~92% · reference_assembly ~68% · assay_type and platform ~8–9%. The ~3% missing even on data_type is mostly a ~16K tail of formats with no rules yet — those files are uncovered on every dimension"),
        ("FASTQ", "platform 100% (read names encode the instrument) — but modality/assay 100% not_classified: a WGS and an RNA-seq FASTQ carry no in-file marker of which they are"),
        ("BED", "reference recoverable only when intervals span whole chromosomes (sparse files are ambiguous); coordinates carry no instrument or experiment signature — assay is known for just 40 of 11,523 BEDs (filename tokens), platform for none; the rest split ~62% not_applicable / ~38% not_classified"),
        ("The lesson", "the hardest dimensions are declared in the submission but invisible in the content — no engine can extract what is not in the bytes"),
    ],
    size=18,
)
add_logo(slide)

# ------------------------------- Slide 14: segue, upcoming sources of truth
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Upcoming: more sources of truth"
color_title(slide, ANVIL_BLUE)
term_bullets(
    slide,
    [
        ("Import project metadata", "read the project's own catalog as ground truth (e.g. the HPRC Catalog's library_strategy — fills the FASTQ modality/assay gap)"),
        ("dbGaP FHIR API", "study-level records constrain the possibilities for every file in a study"),
        ("Marker / methods papers", "find each study's flagship publication and mine its methods section"),
        ("UCSC Genome Browser", "to explore for open-access BED files: a track hub cannot exist without declaring its assembly (genomes.txt) — the exact fact BED content cannot reveal"),
        ("Where else should we look?", "open question for this group — which authorities do you trust and use?"),
    ],
    size=21,
)
add_logo(slide)

# --------------------------------- Slide 17: marker papers — the approach
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Finding a study's marker paper — the approach"
color_title(slide, ANVIL_BLUE)
term_bullets(
    slide,
    [
        ("Marker paper", "the publication that describes the study itself — design, cohort, methods — as opposed to papers that merely use the data"),
        ("Ordered source chain", "ten sources tried in order: dbGaP's own Selected Publications list, the dbGaP study page (PIs, grants, study prose), PMC full-text and PubMed accession searches, dbGaP FHIR citations, NIH RePORTER grant→publication links, the dataset record's leads, cohort-name search, web search last"),
        ("Deterministic fetches + judgment", "a small CLI does the API calls; selection judges titles, years, list position, and grant-link counts — a paper linked by all of a study's grants is a consortium paper"),
        ("Then abstracts are checked", "each candidate's abstract is verified against the study record on identifying facts (cohort size, geography, institution) — this catches wrong candidates; full methods reading is the next epic"),
        ("Candidates, not conclusions", "selection yields a marker-paper candidate per study; each candidate (87 studies) needs manual validation before it is trusted"),
        ("Authoritative path", "in parallel, we have reached out to the dbGaP team about a more authoritative study→publication linkage — others are requesting this capability as well"),
    ],
    size=18,
)
add_logo(slide)

# ----------------------------------- Slide 18: marker-paper candidate results
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Marker-paper candidates — where we stand"
color_title(slide, ANVIL_BLUE)
term_bullets(
    slide,
    [
        ("Coverage", "77 of 87 studies have a marker-paper candidate (67/74 phs-anchored · 10/13 anchor-less)"),
        ("From dbGaP directly", "42/74 studies publish a Selected Publications list — the lead entry is typically the marker for cohort/project studies, but only a candidate for sequencing-center deposits"),
        ("From fallback search", "22 of the 29 studies with empty/absent lists got candidates via grant→publication links and cohort-name search"),
        ("Spot-verified", "candidate abstracts checked against study records on identifying facts — 11 verified; 1 wrong candidate caught and corrected (phs002205 is the Global Microbiome Conservancy, not an IBD cohort)"),
        ("Unresolved", "7 studies with no candidate — mostly center cohorts that appear unpublished as such (WUCADS, NUgene, Cleveland GeneBank, …)"),
        ("Readable", "the 77 studies share candidates — 50 distinct papers in all (consortium papers cover many studies); 42 of the 50 have free PMC full text, so methods sections are one API call away"),
    ],
    size=17,
)
add_logo(slide)

# --------------------------------- Slide: consent-code validation via FHIR
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Validating consent codes against dbGaP FHIR"
color_title(slide, ANVIL_TEAL)
term_bullets(
    slide,
    [
        ("The check", "every AnVIL consent label is validated against its study's registered consents in the dbGaP FHIR record (StudyConsents) — flag-only, nothing is rejected"),
        ("Snapshot", "207 labels across 87 studies: 169 dbGaP-registered · 12 open-access labels · 1 placeholder (TBD) · 7 malformed · 18 unmatched — sums to 207"),
        ("Open access is AnVIL-side", "dbGaP registers no open/unrestricted code at all — NRES (a GA4GH DUO term, 11 datasets) and free-text 'Unrestricted access' (1) are platform labels, not dbGaP vocabulary"),
        ("Gap: mismatches need adjudication", "18 unmatched (dbGaP-style codes the study's registry doesn't list, plus free-text one-offs like 'Consortia Access Only') and 7 underscore-corrupted labels on one study — detection hardening tracked in #328"),
    ],
    size=17,
)
add_logo(slide)

# ------------------------------ Slide: what else methods sections can populate
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "What else the methods sections could populate"
color_title(slide, ANVIL_TEAL)
term_bullets(
    slide,
    [
        ("organism", "stated plainly in methods and study records — human vs non-human primate (dGTEx) vs mouse (IGVF); a dimension we do not classify today"),
        ("tissue / sample type", "methods name the source material — LCLs, blood, brain, tumor; anatomy/tissue modeling is a known follow-on"),
        ("Our five, filled at study level", "methods state assay and platform outright — e.g. NIA CARD: 'ONT long-read WGS at ~30×' — exactly the fields file content cannot reveal"),
        ("Population / cohort facts", "ancestry composition, cohort size, disease focus — study-level context for cohort building"),
    ],
    size=18,
)
add_logo(slide)

# ------------------------------ Slide: which of our five fields are new to AnVIL
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "New and enriched metadata fields"
color_title(slide, ANVIL_BLUE)
term_bullets(
    slide,
    [
        ("Enriching existing fields", "data_modality and reference_assembly already exist per file — submissions rarely populate them, and Meta-Disco can fill them at scale"),
        ("Adding new fields", "data_type, assay_type, and platform would be new per-file dimensions"),
        ("Why it matters", "these fields are exactly what file-level search and cohort building need: what the file contains, what experiment produced it, on which instrument"),
    ],
    size=19,
)
add_logo(slide)

# --------------------------------- Slide: allowed values of the new fields
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = "The new fields' controlled vocabularies"
color_title(slide, ANVIL_BLUE)
add_code_block(
    slide,
    """# from the LinkML schema — every value is enumerated, nothing free-text
data_type (26):
    alignments   reads   sequence   assembly   assembly.reference
    pangenome    pangenome.reference
    variants     variants.germline   variants.somatic   variants.structural   variants.cnv
    genotypes    expression_matrix   quantification   annotations   peaks
    signal       raw_signal   array_signal   images   interval_set
    index        checksum     statistics    log

assay_type (8):
    WGS   WES   RNA-seq   ATAC-seq   ChIP-seq   Bisulfite-seq   Methylation array   Histology

platform (6):
    ILLUMINA   PACBIO   ONT   MGI   ELEMENT   ULTIMA""",
    Inches(0.8),
    Inches(1.5),
    Inches(11.8),
    Inches(5.1),
    size=14,
    spacing=1.05,
    accent_prefixes=("data_type", "assay_type", "platform"),
)
add_logo(slide)

# ------------------------------------ Slide: what we would add to AnVIL today
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "What we would add to AnVIL today"
color_title(slide, ANVIL_BLUE)
term_bullets(
    slide,
    [
        ("~1.85M field values", "committed classifications on files where AnVIL's metadata is blank — or has no such field at all (2026-08 run joined to the AnVIL snapshot)"),
        ("data_modality", "617,404 added — AnVIL's own field is filled on just 6,755 of 734K files (0.9%)"),
        ("reference_assembly", "463,160 added — AnVIL's field filled on only 4,696"),
        ("data_type · assay · platform", "652,265 · 54,349 · 57,973 added — AnVIL has no per-file fields for these at all"),
        ("The reverse is tiny", "AnVIL knows something we don't on ~6,300 modality and ~4,100 reference files — cases to learn from, not a gap"),
    ],
    size=18,
)
add_logo(slide)

# ------------------------------------------ Slide 19: closing, delivery path
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Closing the loop: delivering metadata to AnVIL"
color_title(slide, ANVIL_BLUE)
term_bullets(
    slide,
    [
        ("Set up incremental delivery", "Meta-Disco starts delivering metadata to AnVIL on an ongoing basis — not a one-off backfill"),
        ("Preferred: through TDR", "add the metadata to the TDR input workspaces, so it flows through the existing ingest path and lands where all other metadata lives"),
        ("Alternative: sidecar database", "a separate store carrying the additional metadata plus provenance, which Azul references independently of TDR"),
        ("Provenance travels with the data", "either way, every delivered value keeps its evidence trail — which rule fired, from which source, and why"),
    ],
    size=19,
)
add_logo(slide)

# ---------------------------------------------------- Slide: thank you
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Thank you!"
slide.placeholders[1].text = "Questions?"
color_title(slide, ANVIL_BLUE)
for para in slide.placeholders[1].text_frame.paragraphs:
    para.font.color.rgb = ANVIL_TEAL
    para.font.size = Pt(28)
add_logo(slide)

# ---------------------------------------------------------------- save
brand_fonts(prs)
linkify_issue_refs(prs)
out = sys.argv[1]
prs.save(out)
print("wrote", out, f"({len(prs.slides.__iter__().__length_hint__() if hasattr(prs.slides, '__length_hint__') else list(prs.slides))} slides)" if False else f"({len(prs.slides._sldIdLst)} slides)")
